#!/usr/bin/env python3
"""
Project Introduction PPT Template
Copy and modify this file for new project presentations.

Usage:
    python project-ppt-template.py

Customize:
    - PROJECT_NAME, PROJECT_TAGLINE
    - BRAND_COLORS
    - slides_data structure
    - shot images path
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
import os

# ==================== CONFIGURATION ====================

PROJECT_NAME = "Your Project Name"
PROJECT_TAGLINE = "Project tagline or description"

# Brand colors (customize for your project)
COLORS = {
    "primary": RGBColor(0x10, 0x24, 0x3f),    # Deep blue
    "accent": RGBColor(0xc9, 0xa2, 0x4b),      # Warm gold
    "background_light": RGBColor(0xf3, 0xed, 0xe1),  # Rice white
    "text_dark": RGBColor(0x10, 0x24, 0x3f),
    "text_light": RGBColor(0xf3, 0xed, 0xe1),
}

# ==================== SLIDE CONTENT ====================

def create_title_slide(prs, title, subtitle):
    """Create title slide with dark background."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, 
        prs.slide_width, prs.slide_height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["primary"]
    
    title_box = slide.shapes.add_textbox(
        Inches(2), Inches(2), Inches(9), Inches(2)
    )
    tf = title_box.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(48)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = COLORS["accent"]
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    subtitle_box = slide.shapes.add_textbox(
        Inches(2), Inches(4.5), Inches(9), Inches(1.5)
    )
    tf2 = subtitle_box.text_frame
    tf2.text = subtitle
    tf2.paragraphs[0].font.size = Pt(24)
    tf2.paragraphs[0].font.color.rgb = COLORS["background_light"]
    tf2.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    return slide

def create_problem_solution_slide(prs, problems, solutions):
    """Create problem/solution comparison slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # ... implementation
    return slide

def create_architecture_slide(prs, modules):
    """Create architecture overview slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # ... implementation
    return slide

# ==================== MAIN ====================

def create_presentation():
    """Main function to create the presentation."""
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Title slide
    create_title_slide(prs, PROJECT_NAME, PROJECT_TAGLINE)
    
    # Add more slides here...
    # create_problem_solution_slide(prs, problems, solutions)
    # create_architecture_slide(prs, modules)
    
    # Save
    output_path = f"{PROJECT_NAME.replace(' ', '_')}_Introduction.pptx"
    prs.save(output_path)
    print(f"Saved: {output_path}")
    return output_path

if __name__ == "__main__":
    create_presentation()